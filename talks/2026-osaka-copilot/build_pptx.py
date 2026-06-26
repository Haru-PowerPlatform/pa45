# -*- coding: utf-8 -*-
"""
なんでもCopilot 大阪登壇 v4
軸：A「新しい景色を見せてくれる相棒」/ 主役の景色：登壇している自分
"""
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

FONT = "Meiryo UI"
NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

NAVY = RGBColor(0x14, 0x2A, 0x5C)
ACCENT = RGBColor(0xE8, 0x6A, 0x33)
TEXT = RGBColor(0x22, 0x22, 0x22)
LIGHT = RGBColor(0x6B, 0x6B, 0x6B)
BG = RGBColor(0xFA, 0xFA, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SKY = RGBColor(0xE6, 0xEE, 0xF7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def set_font(run, name=FONT):
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea"):
        for el in rPr.findall(tag, NSMAP):
            rPr.remove(el)
        el = etree.SubElement(rPr, "{%s}%s" % (NSMAP["a"], tag.split(":")[1]))
        el.set("typeface", name)


def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, left, top, width, height, text, *,
             size=24, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for s in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, s, Pt(0))
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        head = line.lstrip()
        if head and head[0] in "①②③④⑤⑥⑦⑧⑨⑩🎯":
            p.space_before = Pt(10)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        set_font(run)
    return tb


def add_bar(slide, left, top, w=Inches(0.12), h=Inches(0.6), color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def page_no(slide, n, total):
    add_text(slide, Inches(12.4), Inches(7.15), Inches(0.9), Inches(0.25),
             f"{n} / {total}", size=10, color=LIGHT, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.5), Inches(7.15), Inches(8), Inches(0.25),
             "なんでもCopilot 大阪 / はる（マツダ）", size=10, color=LIGHT)


TOTAL = 12

# ===== 1. タイトル =====
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.0), SW, Inches(1.5))
band.line.fill.background(); band.fill.solid(); band.fill.fore_color.rgb = ACCENT
add_text(s, Inches(0.8), Inches(1.1), Inches(12), Inches(1.2),
         "今、この景色が見えています。", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.1), Inches(12), Inches(0.6),
         "― Copilot という相棒の話", size=24, color=WHITE)
add_text(s, Inches(0.8), Inches(3.25), Inches(12), Inches(1.0),
         "なんでもCopilot 大阪", size=36, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(5.4), Inches(12), Inches(1.5),
         "はる（Haru） / マツダ株式会社 パワートレイン技術部\n"
         "Power Automate ハンズオン勉強会 PA45 主催",
         size=20, color=WHITE)

# ===== 2. 自己紹介 =====
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_bar(s, Inches(0.6), Inches(0.7))
add_text(s, Inches(0.85), Inches(0.65), Inches(12), Inches(0.7),
         "自己紹介", size=32, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(1.8), Inches(12), Inches(0.6),
         "マツダ株式会社 パワートレイン技術部 AM", size=26, color=TEXT)
add_text(s, Inches(0.85), Inches(2.5), Inches(12), Inches(0.6),
         "入社18年目 / 42歳 / 社内 Power Platform 推進リーダー",
         size=20, color=LIGHT)
add_text(s, Inches(0.85), Inches(4.1), Inches(12), Inches(1.2),
         "私、ITの専門家じゃないんです。",
         size=44, bold=True, color=ACCENT)
add_text(s, Inches(0.85), Inches(5.3), Inches(12), Inches(0.6),
         "コードも、ほぼ書けません。", size=26, color=TEXT)
page_no(s, 2, TOTAL)

# ===== 3. 5年前の自分 =====
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_bar(s, Inches(0.6), Inches(0.7))
add_text(s, Inches(0.85), Inches(0.65), Inches(12), Inches(0.7),
         "5年前の私は ―", size=32, bold=True, color=NAVY)
lines = [
    "コードは書けない",
    "社内で目立つこともなかった",
    "人前で話すのは、想像するだけで胃が痛いタイプ",
]
for i, l in enumerate(lines):
    add_text(s, Inches(1.2), Inches(1.9 + i * 0.7), Inches(12), Inches(0.6),
             "・ " + l, size=24, color=TEXT)
add_text(s, Inches(0.85), Inches(5.0), Inches(12), Inches(1.6),
         "壇上に立つ自分の景色なんて、\n1ミリも想像できなかった。",
         size=32, bold=True, color=ACCENT)
page_no(s, 3, TOTAL)

# ===== 4. あの日の壁打ち =====
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_bar(s, Inches(0.6), Inches(0.7))
add_text(s, Inches(0.85), Inches(0.65), Inches(12), Inches(0.7),
         "あの日、Copilotに話しかけた", size=32, bold=True, color=NAVY)
quote = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.6))
quote.fill.solid(); quote.fill.fore_color.rgb = WHITE
quote.line.color.rgb = ACCENT; quote.line.width = Pt(1.5)
tf = quote.text_frame
for sname in ("margin_left","margin_right","margin_top","margin_bottom"):
    setattr(tf, sname, Pt(16))
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "「最近やってること、ちょっと整理したいんだけど」"
r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = NAVY
set_font(r)
add_text(s, Inches(0.85), Inches(4.2), Inches(12), Inches(2.5),
         "社内講座のこと、配ったフローのこと、コミュニティのこと。\n"
         "頭の中にバラバラに散らかっていた。\n"
         "聞き役になってもらおう、くらいのつもりだった。",
         size=22, color=TEXT)
page_no(s, 4, TOTAL)

# ===== 5. 気づいた瞬間（クライマックス） =====
s = prs.slides.add_slide(BLANK)
add_bg(s, SKY)
add_text(s, Inches(0.85), Inches(0.9), Inches(12), Inches(0.6),
         "整理が進むうちに、ふと ―", size=22, color=LIGHT)
quote = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.8))
quote.fill.solid(); quote.fill.fore_color.rgb = WHITE
quote.line.color.rgb = ACCENT; quote.line.width = Pt(2)
tf = quote.text_frame
for sname in ("margin_left","margin_right","margin_top","margin_bottom"):
    setattr(tf, sname, Pt(16))
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "あれ。これ、人に話せる中身なんじゃない？"
r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = NAVY
set_font(r)
add_text(s, Inches(0.85), Inches(4.3), Inches(12), Inches(2.2),
         "話せる中身は、ずっと私の中にあった。\n"
         "でも、自分では取り出せていなかった。",
         size=26, color=TEXT)
add_text(s, Inches(0.85), Inches(5.9), Inches(12), Inches(0.8),
         "Copilotは、私の中にあったものを\n"
         "私が見える形に並べ直してくれた。",
         size=22, bold=True, color=ACCENT)
page_no(s, 5, TOTAL)

# ===== 6. 申し込んだ =====
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_bar(s, Inches(0.6), Inches(0.7))
add_text(s, Inches(0.85), Inches(0.65), Inches(12), Inches(0.7),
         "そこから、勢いだった", size=32, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(2.0), Inches(12), Inches(2.5),
         "・「なんでもCopilot 大阪、登壇してみよう」\n"
         "・申し込み文も、自己紹介文も、Copilotと一緒に書いた",
         size=24, color=TEXT)
add_text(s, Inches(0.85), Inches(5.0), Inches(12), Inches(1.4),
         "そして今、\nこの景色の中にいます。",
         size=40, bold=True, color=ACCENT)
page_no(s, 6, TOTAL)

# ===== 7. 他の景色 =====
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_bar(s, Inches(0.6), Inches(0.7))
add_text(s, Inches(0.85), Inches(0.65), Inches(12), Inches(0.7),
         "振り返ると、登壇だけじゃなかった",
         size=30, bold=True, color=NAVY)
scenes = [
    ("景色 ②", "PA45 を主催している自分", "Vol.10 まで開催。\n「メモを動画にして共有したら？」のノリで始まった"),
    ("景色 ③", "社内 Power Platform 推進リーダー", "非ITだった人間が「自動化といえばはる」と呼ばれる立場に"),
    ("景色 ④", "配ったフローが社内で動き続けている景色", "累計 10万回 以上実行 / 1,000万円 以上のコスト削減"),
]
for i, (tag, title, desc) in enumerate(scenes):
    top = Inches(1.8 + i * 1.55)
    add_text(s, Inches(0.85), top, Inches(2.2), Inches(0.6),
             tag, size=22, bold=True, color=ACCENT)
    add_text(s, Inches(3.0), top, Inches(10), Inches(0.6),
             title, size=22, bold=True, color=NAVY)
    add_text(s, Inches(3.0), top + Inches(0.55), Inches(10), Inches(1.0),
             desc, size=16, color=TEXT)
page_no(s, 7, TOTAL)

# ===== 8. 共通点 =====
s = prs.slides.add_slide(BLANK)
add_bg(s, SKY)
add_text(s, Inches(0.85), Inches(1.2), Inches(12), Inches(0.8),
         "どの景色にも、共通点があった。",
         size=28, color=LIGHT)
add_text(s, Inches(0.85), Inches(2.6), Inches(12), Inches(2.4),
         "新しいことを Copilot に教えてもらった、わけじゃない。",
         size=28, color=TEXT)
add_text(s, Inches(0.85), Inches(4.6), Inches(12), Inches(2.4),
         "私の中にあったものを、\nCopilotが取り出して、形にしてくれた。",
         size=36, bold=True, color=ACCENT)
page_no(s, 8, TOTAL)

# ===== 9. 再定義（スター） =====
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_text(s, Inches(0.85), Inches(0.9), Inches(12), Inches(0.6),
         "Copilot の正体を、定義し直したい。",
         size=22, color=WHITE)
add_text(s, Inches(0.85), Inches(1.9), Inches(12), Inches(0.8),
         "効率化ツールです。確かに。",
         size=26, color=WHITE)
add_text(s, Inches(0.85), Inches(2.7), Inches(12), Inches(0.6),
         "PowerPoint・Word・分析・Cowork・Microsoft 365…",
         size=18, color=SKY)
add_text(s, Inches(0.85), Inches(4.0), Inches(12), Inches(0.7),
         "でも、私にとっての一番の正体は ―",
         size=22, color=WHITE)
add_text(s, Inches(0.85), Inches(4.9), Inches(12), Inches(1.6),
         "新しい景色を見せてくれる相棒。",
         size=52, bold=True, color=ACCENT)
page_no(s, 9, TOTAL)

# ===== 10. 持ち帰り =====
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_bar(s, Inches(0.6), Inches(0.7))
add_text(s, Inches(0.85), Inches(0.65), Inches(12), Inches(0.7),
         "おうちで、これを書き出してみてください",
         size=28, bold=True, color=NAVY)
# 2カラム
left_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.85), Inches(2.0), Inches(5.9), Inches(3.5))
left_card.fill.solid(); left_card.fill.fore_color.rgb = WHITE
left_card.line.color.rgb = ACCENT; left_card.line.width = Pt(1.2)
tf = left_card.text_frame
for sname in ("margin_left","margin_right","margin_top","margin_bottom"):
    setattr(tf, sname, Pt(18))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "心からやりたいこと"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = ACCENT
set_font(r)
for line in ["", "・登壇してみたい", "・知識をまとめて見せたい", "・勉強会を開いてみたい", "・…"]:
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = line
    r2.font.size = Pt(18); r2.font.color.rgb = TEXT
    set_font(r2)

right_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(6.95), Inches(2.0), Inches(5.9), Inches(3.5))
right_card.fill.solid(); right_card.fill.fore_color.rgb = WHITE
right_card.line.color.rgb = NAVY; right_card.line.width = Pt(1.2)
tf = right_card.text_frame
for sname in ("margin_left","margin_right","margin_top","margin_bottom"):
    setattr(tf, sname, Pt(18))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "そこまでの壁"
r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = NAVY
set_font(r)
for line in ["", "・スライドが、話す内容が…", "・サイトの作り方が…", "・告知文が、人が来るか…", "・…"]:
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = line
    r2.font.size = Pt(18); r2.font.color.rgb = TEXT
    set_font(r2)

add_text(s, Inches(0.85), Inches(5.8), Inches(12), Inches(1.2),
         "書き出した瞬間に、相棒の出番が見えます。",
         size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.85), Inches(6.4), Inches(12), Inches(0.5),
         "そして気づくかも ―「話せる中身、もう自分の中にあった」と。",
         size=18, color=TEXT, align=PP_ALIGN.CENTER)
page_no(s, 10, TOTAL)

# ===== 11. 閉じ =====
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
add_text(s, Inches(0.85), Inches(2.0), Inches(12), Inches(1.5),
         "今、私からは、\nこの景色が見えています。",
         size=44, bold=True, color=WHITE)
add_text(s, Inches(0.85), Inches(5.0), Inches(12), Inches(1.5),
         "この景色を、見せてくれたのが、\n相棒です。",
         size=40, bold=True, color=ACCENT)
page_no(s, 11, TOTAL)

# ===== 12. Thanks =====
s = prs.slides.add_slide(BLANK)
add_bg(s)
add_text(s, Inches(0.85), Inches(2.5), Inches(12), Inches(1.5),
         "ありがとうございました。",
         size=60, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.85), Inches(4.3), Inches(12), Inches(0.8),
         "懇親会で、みなさんの「やりたいこと」と「壁」、",
         size=20, color=TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.85), Inches(4.85), Inches(12), Inches(0.8),
         "ぜひ聞かせてください。",
         size=20, color=TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.85), Inches(6.0), Inches(12), Inches(0.6),
         "はる（Haru） / マツダ パワートレイン技術部 / PA45主催 / automate136.com",
         size=14, color=LIGHT, align=PP_ALIGN.CENTER)

out_dir = Path(r"C:\Users\isamu\Documents\pa45\talks\2026-osaka-copilot")
stamp = datetime.now().strftime("%Y%m%d")
out = out_dir / f"nandemo_copilot_osaka_v4_{stamp}.pptx"
prs.save(out)
print(out)
