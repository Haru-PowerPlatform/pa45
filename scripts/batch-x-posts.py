"""
X投稿スライド一括処理スクリプト
- OGP画像生成
- WordPress下書き作成
- activities JSON作成
- activities-index.json更新

使い方:
  python scripts/batch-x-posts.py
"""

import sys, io, os, json, base64, subprocess, shutil, time
from pathlib import Path
from urllib.request import urlopen, Request
import urllib.error

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

ROOT = Path(__file__).parent.parent
SLIDES_DIR = ROOT / "assets" / "x" / "slides"
ASSETS_X   = ROOT / "assets" / "x"
ACTIVITIES = ROOT / "data" / "activities"
INDEX_PATH = ROOT / "data" / "meta" / "activities-index.json"
SCRIPTS    = ROOT / "scripts"

AVATAR = "https://www.automate136.com/wp-content/uploads/2025/07/u9429395585_side_profile_portrait_of_a_man_in_his_30s_with_me_e4cca787-50fc-42ba-8f64-d0185bce97e5_1.png"

# ── 環境変数 ──────────────────────────────────────────────────────────────────
def load_env():
    env = {}
    with open(ROOT / ".env", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            k, _, v = line.partition("=")
            env[k.strip()] = v.strip()
    return env

ENV = load_env()

def wp(method, endpoint, data=None, binary=None, content_type="application/json", extra_headers=None):
    url   = f"{ENV['WP_URL']}/wp-json/wp/v2/{endpoint}"
    creds = base64.b64encode(f"{ENV['WP_USER']}:{ENV['WP_PASS']}".encode()).decode()
    headers = {"Authorization": f"Basic {creds}", "Content-Type": content_type}
    if extra_headers:
        headers.update(extra_headers)
    body = json.dumps(data).encode("utf-8") if data is not None else binary
    req  = Request(url, data=body, headers=headers, method=method)
    try:
        with urlopen(req) as resp:
            return json.loads(resp.read())
    except urllib.error.HTTPError as e:
        print(f"  [WP ERROR] {e.code}: {e.read().decode()[:200]}")
        return None

# ── スライドデータ定義 ─────────────────────────────────────────────────────────
# (vol, date, title, ogp_title, slug, png_src, category, theme)
SLIDE_DATA = [
    # Copilot Studio Vol.1-10
    (1,  "2026-02-15", "Copilot Studio × Power Automate × Word連携　つまずきポイント3選",
     "Copilot Studio\n× Word連携 つまずき3選", "x-vol-01-cs-word-integration",
     "cs_slide01.png", "Copilot Studio", "green"),
    (2,  "2026-02-17", "Copilot Studio × Power Automateの連携ステップ",
     "Copilot Studio\n× PA連携ステップ", "x-vol-02-cs-pa-steps",
     "cs_slide02.png", "Copilot Studio", "green"),
    (3,  "2026-02-19", "Wordのコンテンツコントロール設定方法（前編）",
     "Wordのコンテンツ\nコントロール設定①", "x-vol-03-word-content-control-1",
     "cs_slide03.png", "Copilot Studio", "green"),
    (4,  "2026-02-21", "Wordのコンテンツコントロール設定方法（後編）",
     "Wordのコンテンツ\nコントロール設定②", "x-vol-04-word-content-control-2",
     "cs_slide04.png", "Copilot Studio", "green"),
    (5,  "2026-02-23", "Copilot Studio｜質問ノードは聞く順番を意識すると会話が自然になる",
     "Copilot Studio\n質問ノードの順番", "x-vol-05-cs-question-order",
     "cs_slide05.png", "Copilot Studio", "green"),
    (6,  "2026-02-25", "Copilot Studio｜Topicsの左右の使い方",
     "Copilot Studio\nTopicsの左右の使い方", "x-vol-06-cs-topics",
     "cs_slide06.png", "Copilot Studio", "green"),
    (7,  "2026-02-27", "Copilot Studio｜ホーム画面4つのアイコンの役割",
     "Copilot Studio\nホーム画面アイコン解説", "x-vol-07-cs-home-icons",
     "cs_slide07.png", "Copilot Studio", "green"),
    (8,  "2026-03-01", "Copilot Studio｜ツール（エージェントフロー）の役割",
     "Copilot Studio\nツールの役割", "x-vol-08-cs-tool-role",
     "cs_slide08.png", "Copilot Studio", "green"),
    (9,  "2026-03-03", "Copilot Studio｜カスタムプロンプト画面の役割",
     "Copilot Studio\nカスタムプロンプト解説", "x-vol-09-cs-custom-prompt",
     "cs_slide09.png", "Copilot Studio", "green"),
    (10, "2026-03-05", "Copilot Studio｜カスタム変数とシステム変数の違い",
     "Copilot Studio\n変数の種類と違い", "x-vol-10-cs-variables",
     "cs_slide10.png", "Copilot Studio", "green"),
    # Vol.11 スキップ（空スライド）
    # Power Automate Vol.12-30
    (12, "2026-03-07", "Power Automate：結合(JOIN)で配列を1本の文字列にまとめる",
     "Power Automate\n結合(JOIN)を使いこなす", "x-vol-12-pa-join",
     "pa_slide01.png", "Power Automate", "blue"),
    (13, "2026-03-09", "Power Automate：1アクションでTeams通知を自動化する",
     "Power Automate\nTeams通知を自動化", "x-vol-13-pa-teams-notify",
     "pa_slide02.png", "Power Automate", "blue"),
    (14, "2026-03-10", "Power Automate：Teams通知・送信1アクションでできること",
     "Power Automate\n送信1アクションの底力", "x-vol-14-pa-teams-send",
     "pa_slide03.png", "Power Automate", "blue"),
    (15, "2026-03-11", "Power Automate：会議の作成アクションで予定を自動化する",
     "Power Automate\n会議の作成アクション", "x-vol-15-pa-create-meeting",
     "pa_slide04.png", "Power Automate", "blue"),
    (16, "2026-03-12", "Power Automate：JSONの解析が苦手な人向け・Composeで軽量処理",
     "Power Automate\nJSON軽量処理の考え方", "x-vol-16-pa-json-compose",
     "pa_slide05.png", "Power Automate", "blue"),
    (17, "2026-03-13", "Power Automate：【作成】アクションでライトなJSON処理をする",
     "Power Automate\n【作成】でJSON処理", "x-vol-17-pa-json-create",
     "pa_slide06.png", "Power Automate", "blue"),
    (18, "2026-03-14", "Power Automate：JSONの構造をイラストで理解する",
     "Power Automate\nJSONをイラストで理解", "x-vol-18-pa-json-image",
     "pa_slide07.png", "Power Automate", "blue"),
    (19, "2026-03-15", "Power Automate：【JSONの解析】アクションで動的コンテンツを使う",
     "Power Automate\nJSONの解析アクション", "x-vol-19-pa-parse-json",
     "pa_slide08.png", "Power Automate", "blue"),
    (20, "2026-03-16", "Power Automate：【フィルターアレイ】で必要なデータだけを残す",
     "Power Automate\nフィルターアレイ解説", "x-vol-20-pa-filter-array",
     "pa_slide09.png", "Power Automate", "blue"),
    (21, "2026-03-17", "Power Automate：ハマりがちな7つの問題と解決方法",
     "Power Automate\nハマる7選と解決策", "x-vol-21-pa-trouble-7",
     "pa_slide11.png", "Power Automate", "blue"),
    (22, "2026-03-18", "Power Automate：条件(Condition)のequalsとcontainsを正しく使い分ける",
     "Power Automate\n条件の勘違いを解消", "x-vol-22-pa-condition-equals",
     "pa_slide12.png", "Power Automate", "blue"),
    (23, "2026-03-19", "Power Automate：実務で効く通知設計・Before→After",
     "Power Automate\n実務で効く通知設計", "x-vol-23-pa-notify-design",
     "pa_slide14.png", "Power Automate", "blue"),
    (24, "2026-03-20", "Power Automate：Condition（条件）は「質問をつくるアクション」",
     "Power Automate\nConditionは質問アクション", "x-vol-24-pa-condition-question",
     "pa_slide15.png", "Power Automate", "blue"),
    (25, "2026-03-21", "Power Automate：翌月1日を毎月自動で出す方法",
     "Power Automate\n翌月1日を自動で出す", "x-vol-25-pa-next-month-1st",
     "pa_slide16.png", "Power Automate", "blue"),
    (26, "2026-03-22", "Power Automate：前月末は「月初から1日戻る」で一発",
     "Power Automate\n前月末を確実に出す", "x-vol-26-pa-prev-month-end",
     "pa_slide17.png", "Power Automate", "blue"),
    # Vol.27 (x-last-bizday) はスキップ（既存活動として登録済み）
    (28, "2026-03-23", "Power Automate：毎月の「◯営業日目」を自動で出す方法",
     "Power Automate\n毎月の◯営業日を算出", "x-vol-28-pa-nth-bizday",
     "pa_slide19.png", "Power Automate", "blue"),
    (29, "2026-03-24", "Power Automate：if()で条件分岐をシンプルに書く",
     "Power Automate\nif()の使いどころ", "x-vol-29-pa-if-function",
     "pa_slide20.png", "Power Automate", "blue"),
    (30, "2026-03-25", "Power Automate：空欄のせいでフローが止まる問題を解決する",
     "Power Automate\n空欄エラーを撃退", "x-vol-30-pa-empty-fix",
     "pa_slide21.png", "Power Automate", "blue"),
]

# ── スライドテキスト取得 ────────────────────────────────────────────────────────
from pptx import Presentation

def get_slide_texts():
    texts = {}
    files = {
        "cs": ROOT / "assets" / "x" / "CopilotStudioVol1~Vol11.pptx",
        "pa": ROOT / "assets" / "x" / "Update_pa45_X投稿_MasterVol12-30.pptx",
    }
    for prefix, path in files.items():
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, 1):
            key = f"{prefix}_slide{i:02d}.png"
            all_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t and len(t) > 2:
                            all_texts.append(t)
            texts[key] = all_texts
    return texts

# ── ブログ本文生成 ─────────────────────────────────────────────────────────────
def build_blog(vol, title, category, slide_texts):
    def bubble(text):
        return f"""<!-- wp:html -->
<div class="speech-wrap sb-id-14 sbs-stn sbp-l sbis-cb cf">
<div class="speech-person">
<figure class="speech-icon"><img class="speech-icon-image" src="{AVATAR}" alt="" width="1024" height="1024" /></figure>
</div>
<div class="speech-balloon">
{text}
</div>
</div>
<!-- /wp:html -->"""

    # スライドテキストからポイントを抽出
    points = [t for t in slide_texts if len(t) > 5 and not t.startswith("①") and not t.startswith("✅") and not t.startswith("Vol.")][:8]
    check_points = [t for t in slide_texts if t.startswith("✅")][:5]

    content = f"""<!-- wp:paragraph -->
<p>今回のテーマは「<span style="color: #0000ff;"><strong>{title}</strong></span>」です。</p>
<!-- /wp:paragraph -->

<!-- wp:paragraph -->
<p>X（旧Twitter）での1枚スライドをもとに、もう少し詳しく解説します。</p>
<!-- /wp:paragraph -->

&nbsp;

<!-- wp:heading -->
<h2 class="wp-block-heading">ポイント</h2>
<!-- /wp:heading -->
"""

    if check_points:
        content += "<!-- wp:list -->\n<ul class=\"wp-block-list\">"
        for p in check_points:
            content += f"<li>{p}</li>"
        content += "</ul>\n<!-- /wp:list -->\n\n"
    elif points:
        content += "<!-- wp:list -->\n<ul class=\"wp-block-list\">"
        for p in points[:5]:
            content += f"<li>{p}</li>"
        content += "</ul>\n<!-- /wp:list -->\n\n"

    content += f"""
&nbsp;

<!-- wp:heading -->
<h2 class="wp-block-heading">スライドで解説</h2>
<!-- /wp:heading -->

<!-- wp:paragraph -->
<p>以下のスライドにまとめています。↓</p>
<!-- /wp:paragraph -->

<!-- wp:paragraph -->
<p>（スライド画像）</p>
<!-- /wp:paragraph -->

&nbsp;

<!-- wp:heading -->
<h2 class="wp-block-heading">まとめ</h2>
<!-- /wp:heading -->

{bubble(f"Vol.{vol}「{title}」でした。{category}を使う上で知っておくと確実に役立つポイントです。")}
"""
    return content

# ── OGP生成 ──────────────────────────────────────────────────────────────────
def make_ogp(ogp_title, label, theme, out_path):
    result = subprocess.run(
        [sys.executable, str(SCRIPTS / "make-ogp.py"),
         "--title", ogp_title,
         "--label", label,
         "--sub", "1枚でわかる",
         "--theme", theme,
         "--out", str(out_path)],
        cwd=str(ROOT), capture_output=True, text=True
    )
    return result.returncode == 0

# ── WordPress メディアアップロード ─────────────────────────────────────────────
def upload_media(path, filename):
    with open(path, "rb") as f:
        binary = f.read()
    creds = base64.b64encode(f"{ENV['WP_USER']}:{ENV['WP_PASS']}".encode()).decode()
    req = Request(
        f"{ENV['WP_URL']}/wp-json/wp/v2/media",
        data=binary,
        headers={
            "Authorization": f"Basic {creds}",
            "Content-Type": "image/png",
            "Content-Disposition": f"attachment; filename={filename}",
        },
        method="POST"
    )
    try:
        with urlopen(req) as resp:
            return json.loads(resp.read())["id"]
    except Exception as e:
        print(f"  [MEDIA ERROR] {e}")
        return None

# ── activities-index 更新 ──────────────────────────────────────────────────────
def update_index(rel_path):
    if INDEX_PATH.exists():
        with open(INDEX_PATH, encoding="utf-8") as f:
            index = json.load(f)
    else:
        index = []
    if rel_path not in index:
        index.append(rel_path)
        index.sort()
    with open(INDEX_PATH, "w", encoding="utf-8") as f:
        json.dump(index, f, ensure_ascii=False, indent=2)
        f.write("\n")

# ── メイン処理 ────────────────────────────────────────────────────────────────
def main():
    print("スライドテキストを読み込み中...")
    slide_texts = get_slide_texts()

    ogp_tmp = ROOT / "scripts" / "ogp-output.png"

    for vol, date, title, ogp_title, slug, png_src, category, theme in SLIDE_DATA:
        print(f"\n{'='*50}")
        print(f"Vol.{vol:02d}: {title}")

        activity_id = f"{date}-x-vol{vol:02d}"
        activity_file = ACTIVITIES / f"{activity_id}.json"

        # 既存チェック
        if activity_file.exists():
            print(f"  スキップ（既存: {activity_file.name}）")
            continue

        # 1. OGP生成
        print(f"  ▶ OGP生成...")
        label = f"{'Copilot Studio' if category == 'Copilot Studio' else 'Power Automate'} Tips"
        ok = make_ogp(ogp_title, label, theme, ogp_tmp)
        if not ok:
            print(f"  [SKIP] OGP生成失敗")
            continue

        # 2. OGPをWordPressにアップロード
        print(f"  ▶ OGPアップロード...")
        ogp_fname = f"x-vol{vol:02d}-ogp-{int(time.time())}.png"
        media_id = upload_media(ogp_tmp, ogp_fname)
        if not media_id:
            print(f"  [SKIP] アップロード失敗")
            continue
        print(f"  メディアID: {media_id}")

        # 3. WordPress投稿作成
        print(f"  ▶ WordPress下書き作成...")
        texts = slide_texts.get(png_src, [])
        content = build_blog(vol, title, category, texts)
        post_payload = {
            "title": f"【Vol.{vol}】{title}",
            "content": content,
            "status": "draft",
            "slug": slug,
            "featured_media": media_id,
        }
        result = wp("POST", "posts", data=post_payload)
        if not result:
            print(f"  [SKIP] 投稿作成失敗")
            continue
        post_id = result["id"]
        blog_url = result.get("link", "")
        print(f"  投稿ID: {post_id}")

        # 4. スライドPNGをassets/x/にコピー
        src_png = SLIDES_DIR / png_src
        dst_png = ASSETS_X / f"{activity_id}.png"
        if src_png.exists():
            shutil.copy2(src_png, dst_png)

        # 5. activities JSON作成
        activity = {
            "id": activity_id,
            "type": "X",
            "title": f"Vol.{vol}：{title}",
            "date": date,
            "public": True,
            "summary": title,
            "tags": [category, "1枚スライド", "初心者"],
            "evidence": {
                "blog": blog_url,
                "slide": str(dst_png.relative_to(ROOT)).replace("\\", "/"),
            }
        }
        with open(activity_file, "w", encoding="utf-8") as f:
            json.dump(activity, f, ensure_ascii=False, indent=2)
            f.write("\n")
        print(f"  activity JSON: {activity_file.name}")

        # 6. index更新
        update_index(f"data/activities/{activity_file.name}")

        time.sleep(0.5)  # WordPress APIへの負荷軽減

    print(f"\n{'='*50}")
    print("全処理完了！")
    print("次のステップ: git add . && git commit && git push")

if __name__ == "__main__":
    main()
