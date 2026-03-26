"""
PA45 PPTX フォントサイズ修正スクリプト
16pt未満のテキストをすべて16ptに統一する
コピー・.bakファイルはスキップ
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy

ROOT = Path(__file__).parent.parent
PPTX_DIR = ROOT / "assets" / "pa45"
MIN_PT = 16

def fix_pptx(path):
    prs = Presentation(path)
    fixed = 0

    for slide_i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    sz = run.font.size
                    if sz is not None and sz < Pt(MIN_PT):
                        old_pt = sz / 12700
                        run.font.size = Pt(MIN_PT)
                        fixed += 1
                        print(f"  Slide {slide_i} | {old_pt:.1f}pt -> {MIN_PT}pt")
                # paragraph-level font check
                pf = para._pPr
                if pf is not None:
                    for r in para._p.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
                        rpr = r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                        if rpr is not None:
                            sz_attr = rpr.get('sz')
                            if sz_attr and int(sz_attr) < MIN_PT * 100:
                                old_pt = int(sz_attr) / 100
                                rpr.set('sz', str(MIN_PT * 100))
                                fixed += 1

    if fixed:
        prs.save(path)
        print(f"  → 保存完了（{fixed}箇所修正）\n")
    else:
        print(f"  → 修正不要\n")
    return fixed

def main():
    files = sorted([
        f for f in PPTX_DIR.glob("*.pptx")
        if "コピー" not in f.name and ".bak" not in f.name
    ])
    print(f"対象ファイル: {len(files)}件\n")
    total = 0
    for f in files:
        print(f"【{f.name}】")
        total += fix_pptx(f)
    print(f"\n完了: 合計 {total} 箇所を {MIN_PT}pt に修正しました")

if __name__ == "__main__":
    main()
