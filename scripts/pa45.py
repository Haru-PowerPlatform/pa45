"""
PA45 統合コマンド
------------------------------------------------------------
講座前のスライド生成と、講座後の公開を一元管理する。

【講座前】スライドを作る:
  python scripts/pa45.py gen \
    --vol 2 \
    --title-ja "変数を操作しよう" \
    --title-en "Set variable" \
    --date "2026-04-02" \
    --next-vol 3 \
    --next-title-ja "条件分岐を使おう" \
    --next-title-en "Condition" \
    --next-desc "フローに「もし〜なら」の分岐を追加する方法を学びます。"

【講座後】ブログ公開・サイト更新:
  python scripts/pa45.py publish \
    --pptx "outputs/pa45/P002_PA45_変数を操作しよう_20260402.pptx" \
    --vol 2 \
    --title-ja "変数を操作しよう" \
    --title-en "Set variable" \
    --date "2026-04-02" \
    --publish

【確認】スライド構造を解析:
  python scripts/pa45.py inspect \
    --pptx "path/to/file.pptx"
------------------------------------------------------------
"""

import sys
import subprocess
from pathlib import Path

SCRIPT_DIR = Path(__file__).parent


def cmd_gen(argv):
    subprocess.run([sys.executable, str(SCRIPT_DIR / "pa45-gen.py")] + argv)


def cmd_publish(argv):
    subprocess.run([sys.executable, str(SCRIPT_DIR / "pa45-publish.py")] + argv)


def cmd_inspect(argv):
    """PPTXの構造を簡易表示。"""
    import argparse
    p = argparse.ArgumentParser()
    p.add_argument("--pptx", required=True)
    args = p.parse_args(argv)

    try:
        from pptx import Presentation
    except ImportError:
        print("pip install python-pptx")
        sys.exit(1)

    import io, sys
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

    prs = Presentation(args.pptx)
    for i, slide in enumerate(prs.slides):
        print(f"\n=== Slide {i+1} ===")
        for j, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                txt = shape.text_frame.text.replace("\n", "/")
                if txt.strip():
                    print(f"  [{j:2d}] {txt[:80]!r}")


def main():
    if len(sys.argv) < 2 or sys.argv[1] not in ("gen", "publish", "inspect"):
        print(__doc__)
        sys.exit(0)

    cmd  = sys.argv[1]
    rest = sys.argv[2:]

    if cmd == "gen":
        cmd_gen(rest)
    elif cmd == "publish":
        cmd_publish(rest)
    elif cmd == "inspect":
        cmd_inspect(rest)


if __name__ == "__main__":
    main()
