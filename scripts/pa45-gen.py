"""
PA45 スライド自動生成スクリプト
------------------------------------------------------------
テンプレートPPTXをコピーして可変スライドのテキストを差し替える。

使い方:
  python scripts/pa45-gen.py \
    --vol       2 \
    --title-ja  "変数を操作しよう" \
    --title-en  "Set variable" \
    --date      "2026-04-02" \
    --next-vol  3 \
    --next-title-ja "条件分岐を使おう" \
    --next-title-en "Condition" \
    --next-desc "フローに「もし〜なら」の分岐を追加する方法を学びます。" \
    [--template path/to/template.pptx]   # 省略時は最新のPPTX

固定スライド（変更なし）: 1(コンセプト), 2(ルール), 3(Haruプロフ), 4(PA紹介), 6(学習フロー), 16(アンケート)
可変スライド（自動更新）: 5(タイトル), 7(アジェンダ), 15(次回予告)
  + 「PA45｜第N回」「Vol.N」ラベルをすべてのスライドで更新
------------------------------------------------------------
"""

import argparse
import re
import shutil
import sys
from copy import deepcopy
from pathlib import Path

SCRIPT_DIR = Path(__file__).parent
ROOT_DIR   = SCRIPT_DIR.parent
OUTPUT_DIR = ROOT_DIR / "outputs" / "pa45"


# ── python-pptx インポート確認 ─────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("[ERROR] python-pptx が見つかりません: pip install python-pptx")
    sys.exit(1)


# ── テキスト置換ユーティリティ ─────────────────────────────────────────
def _replace_in_paragraph(para, old: str, new: str) -> bool:
    """
    段落内のテキストを置換する。
    テキストが複数 run にまたがる場合は最初の run にまとめる。
    フォーマット（太字・色など）は最初の run を継承。
    """
    full = "".join(r.text for r in para.runs)
    if old not in full:
        return False
    replaced = full.replace(old, new)
    # 最初の run にまとめて書き込む
    if para.runs:
        para.runs[0].text = replaced
        for r in para.runs[1:]:
            r.text = ""
    return True


def replace_text(prs: "Presentation", old: str, new: str, limit: int = 0) -> int:
    """全スライドの全シェイプで old → new を置換。置換件数を返す。"""
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if _replace_in_paragraph(para, old, new):
                    count += 1
                    if limit and count >= limit:
                        return count
    return count


def replace_slide_text(slide, old: str, new: str) -> int:
    """指定スライドのみで置換。"""
    count = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if _replace_in_paragraph(para, old, new):
                count += 1
    return count


def clear_shape_runs(shape):
    """シェイプ内のすべての run のテキストを空にする（段落・書式は保持）。"""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = ""


def set_shape_text(slide, shape_idx: int, new_text: str):
    """
    shape_idx 番目のシェイプのテキストを丸ごと上書き。
    全 run をクリアしてから最初の run に書き込む。
    """
    shape = slide.shapes[shape_idx]
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # 全 run をクリア
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    # 最初の段落・最初の run に書き込む
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = new_text


def replace_shape_all_text(shape, new_text: str):
    """
    シェイプ全体（全段落・全 run）を一括クリアして new_text を書き込む。
    複数行は最初の段落の最初の run にまとめて入る（改行は / で表現できない）。
    Slide 15 のように複数段落にわたる場合に使う。
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # 全 run をクリア
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    # 最初の段落・最初の run に新しいテキストを書き込む
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = new_text
        return
    # run がない場合はスキップ（空段落）
    print(f"  [WARN] shape に run が見つかりません: {shape.name}")


# ── テンプレートを探す ─────────────────────────────────────────────────
def find_template() -> Path:
    candidates = sorted(
        list(ROOT_DIR.glob("**/*PA45*.pptx")) +
        list(Path("C:/Users/isamu/Downloads").glob("**/*PA45*.pptx")),
        key=lambda p: p.stat().st_mtime,
        reverse=True,
    )
    if not candidates:
        print("[ERROR] テンプレートPPTXが見つかりません。--template で指定してください。")
        sys.exit(1)
    return candidates[0]


# ── メイン ────────────────────────────────────────────────────────────
def main():
    p = argparse.ArgumentParser(description="PA45 スライド自動生成")
    p.add_argument("--vol",          required=True, type=int,  help="今回の回数 (例: 2)")
    p.add_argument("--title-ja",     required=True,            help="今回の日本語テーマ (例: 変数を操作しよう)")
    p.add_argument("--title-en",     required=True,            help="今回の英語タイトル (例: Set variable)")
    p.add_argument("--date",         required=True,            help="開催日 YYYY-MM-DD")
    p.add_argument("--next-vol",     required=True, type=int,  help="次回の回数")
    p.add_argument("--next-title-ja",required=True,            help="次回の日本語テーマ")
    p.add_argument("--next-title-en",required=True,            help="次回の英語タイトル")
    p.add_argument("--next-desc",    default="",               help="次回の1行説明")
    p.add_argument("--agenda1-title",default="",               help="アジェンダ1タイトル (\\nで改行)")
    p.add_argument("--agenda1-desc", default="",               help="アジェンダ1説明")
    p.add_argument("--agenda2-title",default="",               help="アジェンダ2タイトル")
    p.add_argument("--agenda2-desc", default="",               help="アジェンダ2説明")
    p.add_argument("--agenda3-title",default="",               help="アジェンダ3タイトル")
    p.add_argument("--agenda3-desc", default="",               help="アジェンダ3説明")
    p.add_argument("--template",     default="",               help="テンプレートPPTXパス")
    args = p.parse_args()

    vol      = args.vol
    prev_vol = vol - 1  # テンプレートの回数（通常 1 つ前）

    # テンプレートPPTX
    template_path = Path(args.template) if args.template else find_template()
    print(f"テンプレート: {template_path}")

    # 出力先
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    safe_title = re.sub(r"[\\/:*?\"<>|]", "_", args.title_ja)
    out_path = OUTPUT_DIR / f"P{vol:03d}_PA45_{safe_title}_{args.date.replace('-','')}.pptx"

    shutil.copy2(template_path, out_path)
    print(f"コピー → {out_path}")

    prs = Presentation(str(out_path))
    slides = prs.slides

    # ── 全スライド共通：回数ラベル・Vol 番号 ──────────────────────────
    prev_label = f"第{prev_vol}回"
    new_label  = f"第{vol}回"
    n = replace_text(prs, prev_label, new_label)
    print(f"  「{prev_label}」→「{new_label}」: {n}箇所")

    n = replace_text(prs, f"Vol.{prev_vol}", f"Vol.{vol}")
    print(f"  「Vol.{prev_vol}」→「Vol.{vol}」: {n}箇所")

    # 次回 Vol
    next_vol_old = f"Vol.{vol}"    # テンプレートの「次回」はちょうど今回の次
    next_vol_new = f"Vol.{args.next_vol}"
    # Slide 15 のみ対象（次回予告スライド = index 14）
    n = replace_slide_text(slides[14], next_vol_old, next_vol_new)
    print(f"  次回Vol: {n}箇所（Slide 15）")

    # ── Slide 5 (index 4)：タイトルスライド ──────────────────────────
    sl5 = slides[4]

    # shape[11]: '第1回：変数を作ろう' → '第N回：テーマ'
    replace_slide_text(sl5, f"第{prev_vol}回：", f"第{vol}回：")

    # タイトルの日本語部分を差し替え（元テキストの後半 "：〇〇" を置換）
    # 元: '第1回：変数を作ろう'  新: '第2回：変数を操作しよう'
    for shape in sl5.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        if "：" in txt and f"第{vol}回：" in txt:
            # "：" の後ろがテーマ名
            for para in shape.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                if f"第{vol}回：" in full:
                    # 末尾のテーマ部分を置換
                    import re as _re
                    new_full = _re.sub(
                        rf"(第{vol}回：).*",
                        rf"\g<1>{args.title_ja}",
                        full,
                    )
                    if new_full != full:
                        _replace_in_paragraph(para, full, new_full)

    # shape[12]: 英語タイトル
    for shape in sl5.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        # 前回の英語タイトルが含まれているシェイプ = "Initialize variable" 的な行
        # パターン: "英語タイトル (日本語)" 形式
        if "(" in txt and ")" in txt and txt == txt.strip():
            # シェイプ全体が1行の英語タイトルと判定
            for para in shape.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                if full.strip() and "(" in full and ")" in full:
                    new_en = f"{args.title_en} ({args.title_ja})"
                    _replace_in_paragraph(para, full, new_en)
                    print(f"  英語タイトル: {full!r} → {new_en!r}")
                    break

    # ── Slide 7 (index 6)：アジェンダ ───────────────────────────────
    sl7 = slides[6]
    if args.agenda1_title:
        agenda_map = [
            (21, 22, args.agenda1_title.replace("\\n", "\n"), args.agenda1_desc),
            (26, 27, args.agenda2_title.replace("\\n", "\n"), args.agenda2_desc),
            (31, 32, args.agenda3_title.replace("\\n", "\n"), args.agenda3_desc),
        ]
        for title_idx, desc_idx, new_title, new_desc in agenda_map:
            if not new_title:
                continue
            # タイトルシェイプ
            shape_t = sl7.shapes[title_idx]
            old_t = shape_t.text_frame.text
            for para in shape_t.text_frame.paragraphs:
                full = "".join(r.text for r in para.runs)
                if full.strip():
                    _replace_in_paragraph(para, full, new_title)
                    print(f"  アジェンダタイトル: {old_t!r} → {new_title!r}")
                    break
            # 説明シェイプ
            if new_desc:
                shape_d = sl7.shapes[desc_idx]
                old_d = shape_d.text_frame.text
                for para in shape_d.text_frame.paragraphs:
                    full = "".join(r.text for r in para.runs)
                    if full.strip():
                        _replace_in_paragraph(para, full, new_desc)
                        print(f"  アジェンダ説明: {old_d[:40]!r} → {new_desc[:40]!r}")
                        break

    # ── Slide 15 (index 14)：次回予告 ────────────────────────────────
    sl15 = slides[14]

    # shape[9]: 次回テーマタイトル（複数段落にまたがる）
    shape_next_title = sl15.shapes[9]
    old_next = shape_next_title.text_frame.text.replace("\n", "/")
    new_next_text = f"{args.next_title_ja}（{args.next_title_en}）"
    replace_shape_all_text(shape_next_title, new_next_text)
    print(f"  次回テーマ: {old_next!r} → {new_next_text!r}")

    # shape[10]: 次回説明文（複数段落）
    if args.next_desc:
        shape_next_desc = sl15.shapes[10]
        old_desc = shape_next_desc.text_frame.text[:40]
        replace_shape_all_text(shape_next_desc, args.next_desc)
        print(f"  次回説明: {old_desc!r} → {args.next_desc[:40]!r}")

    # Slide 15 の「Next Session（第N回）」ラベルも更新
    # テンプレートの「次回」は「第{vol}回」（現在の vol の次を示している）
    # → 新しい PPTX では「第{next_vol}回」に変える
    replace_slide_text(sl15, f"第{vol}回）", f"第{args.next_vol}回）")

    # ── 保存 ─────────────────────────────────────────────────────────
    prs.save(str(out_path))
    print()
    print("=" * 55)
    print(f"[完了] 出力ファイル: {out_path}")
    print(f"[完了] 第{vol}回：{args.title_ja}")
    print(f"[完了] 次回（第{args.next_vol}回）：{args.next_title_ja}")
    print("=" * 55)
    print()
    print("次のステップ（講座後）:")
    print(f"  python scripts/pa45-publish.py \\")
    print(f"    --pptx \"{out_path}\" \\")
    print(f"    --vol {vol} \\")
    print(f"    --title-ja \"{args.title_ja}\" \\")
    print(f"    --date {args.date} \\")
    print(f"    --publish")
    return out_path


if __name__ == "__main__":
    main()
