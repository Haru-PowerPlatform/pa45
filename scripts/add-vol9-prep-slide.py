"""
PA45 第9回 PPTX に「事前準備」スライドを slide 2（既存「20:15スタート」の直後）に挿入。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from copy import deepcopy
from lxml import etree

PPTX = Path('assets/pa45/P009_PA45_SharePointTeams_20260507.pptx')

DOWNLOAD_URL = 'https://haru-powerplatform.github.io/pa45/flows/vol-09/approval-sample.xlsx'

prs = Presentation(str(PPTX))

# Blank layout
blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
slide = prs.slides.add_slide(blank_layout)

# 背景：薄いブルー
from pptx.oxml.ns import qn
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0xF5, 0xFA, 0xFF)

sw = prs.slide_width
sh = prs.slide_height

# タイトルバー
title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), sw, Inches(1.2))
title_bar.fill.solid()
title_bar.fill.fore_color.rgb = RGBColor(0x2A, 0x7D, 0xD4)
title_bar.line.fill.background()
tf = title_bar.text_frame
tf.margin_left = Inches(0.6)
tf.margin_top = Inches(0.25)
p = tf.paragraphs[0]
r = p.add_run()
r.text = '🛠️ ハンズオン 事前準備のお願い（3分でできます）'
r.font.size = Pt(28)
r.font.bold = True
r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# サブタイトル
sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), sw - Inches(1.2), Inches(0.6))
p = sub.text_frame.paragraphs[0]
r = p.add_run()
r.text = '当日のハンズオンで使う「承認用サンプルExcel」を、事前にご自身のSharePointに配置してください'
r.font.size = Pt(16)
r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# ステップ1
step1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.3), sw - Inches(1.2), Inches(1.6))
step1.fill.solid()
step1.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
step1.line.color.rgb = RGBColor(0x2A, 0x7D, 0xD4)
step1.line.width = Pt(1.5)
tf = step1.text_frame
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = 'STEP 1   サンプルExcelをダウンロード'
r.font.size = Pt(20); r.font.bold = True
r.font.color.rgb = RGBColor(0x2A, 0x7D, 0xD4)
p2 = tf.add_paragraph()
p2.space_before = Pt(8)
r = p2.add_run()
r.text = '👇 下のURLをブラウザに貼り付けてください'
r.font.size = Pt(13)
r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
p3 = tf.add_paragraph()
p3.space_before = Pt(4)
r = p3.add_run()
r.text = DOWNLOAD_URL
r.font.size = Pt(15); r.font.bold = True
r.font.color.rgb = RGBColor(0x1A, 0x6B, 0xBF)

# ステップ2
step2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.1), sw - Inches(1.2), Inches(2.3))
step2.fill.solid()
step2.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
step2.line.color.rgb = RGBColor(0x2A, 0x7D, 0xD4)
step2.line.width = Pt(1.5)
tf = step2.text_frame
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = 'STEP 2   ご自身のSharePointにアップロード'
r.font.size = Pt(20); r.font.bold = True
r.font.color.rgb = RGBColor(0x2A, 0x7D, 0xD4)

steps_b = [
    '①  編集権限のあるSharePointサイトを開く（社内サイト・チームサイト等）',
    '②  「ドキュメント」ライブラリを開く',
    '③  「アップロード → ファイル」で approval-sample.xlsx を選択',
    '④  完了！  当日はこのファイルを上司役で更新する想定でフローを作ります',
]
for line in steps_b:
    pp = tf.add_paragraph()
    pp.space_before = Pt(4)
    r = pp.add_run()
    r.text = line
    r.font.size = Pt(14)
    r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# 注意書き
note = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), sw - Inches(1.2), Inches(0.5))
p = note.text_frame.paragraphs[0]
r = p.add_run()
r.text = '※ 配置済みだとハンズオンに専念できます。当日その場で配置してもOK！'
r.font.size = Pt(12); r.font.italic = True
r.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

# 新規追加スライドを末尾→2番目（index 1）に移動
xml_slides = prs.slides._sldIdLst
slides_list = list(xml_slides)
new_slide_xml = slides_list[-1]
xml_slides.remove(new_slide_xml)
xml_slides.insert(1, new_slide_xml)

prs.save(str(PPTX))
print('✓ 事前準備スライドを slide 2 に挿入しました')
print(f'  保存: {PPTX}')
print(f'  全スライド数: {len(prs.slides)}')
