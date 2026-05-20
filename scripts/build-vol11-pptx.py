# -*- coding: utf-8 -*-
"""
PA45 Vol.11 PPTX ビルドスクリプト（v2）
P010（Vol.10）をベースに Vol.11「重要メール仕分けフロー」のデッキを生成する。

処理:
  1. P010 を P011 にコピー
  2. 各スライドの文言を Vol.11 内容へ置換（部分一致キーで shape を特定）
  3. 内容スライドの見出しフォントを 30pt に正規化
  4. Vol.10 のフロー図（大きい画像）を除去 → スクショ貼り付け用プレースホルダを設置
  5. S10/S11 は Vol.10 の事前準備①②を「入力・出力の説明」スライドへ転用

出力: assets/pa45/P011_PA45_RunHistory_20260521.pptx
"""
import sys, io, shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

ROOT = Path(__file__).parent.parent
SRC  = ROOT / "assets" / "pa45" / "P010_PA45_BizReview_20260514.pptx"
DST  = ROOT / "assets" / "pa45" / "P011_PA45_RunHistory_20260521.pptx"

# ── 文言置換: slide番号 -> [(キー部分文字列, 新テキスト全文), ...] ──
REPLACE = {
    1: [
        ("第10回 / Vol.10", "第11回 / Vol.11"),
        ("4つを1本につなぐ総復習", "第11回｜実行履歴を読めると自動化が楽しくなる"),
        ("経費申請の自動化（Forms", "失敗しないフロー設計｜メールの重要度で振り分け"),
        ("Power Automate と Teams を",
         "Power Automate と Outlook を\n事前に起動しておくとスムーズ"),
    ],
    2: [
        ("ハンズオン 事前準備", "事前準備（今回はほぼ不要です）"),
        ("今回は総復習回。事前にSharePoint",
         "今回のハンズオンは特別な準備が要りません。下の2つだけ確認しておきましょう。"),
        ("STEP 1",
         "STEP 1   Power Automate にサインインできる\n👉 make.powerautomate.com を開いて、サインインできる状態にしておきます"),
        ("STEP 2",
         "STEP 2   自分のメール（Outlook）が使える\n👉 ハンズオン中、自分から自分へ1通メールを送ってフローを動かします。それだけです"),
    ],
    3: [
        ("Vol.10", "Vol.11"),
        ("PA45【Forms × SharePoint", "PA45【メール受信 × 条件分岐 × 実行履歴】"),
        ("4つを1本につなぐ総復習",
         "第11回： \n失敗しないフロー設計\n実行履歴を読めるようになろう"),
        ("Theme　Forms", "Theme　メールの重要度で条件分岐｜実行履歴の出力を読む"),
        ("Target　Vol.1〜9",
         "Target　条件分岐（Vol.3）を触った方／はじめての方も、見るだけでOK"),
    ],
    7: [
        ("PA45｜第10回", "PA45｜第11回"),
        ("Formsの回答をSharePointに自動登録する", "メール受信でフローを作る"),
        ("Forms「新しい応答」をトリガーに",
         "「新しいメールが届いたとき」をトリガーに、\n届いたメールで動くフローを作ります。"),
        ("承認フローで上司の判断を待つ", "重要度で条件分岐させる"),
        ("「開始して承認を待機」で承認依頼を送信",
         "「条件」アクションで、メールの重要度が\n「高」かどうかで処理を分けます。"),
        ("結果でSharePoint更新", "実行履歴の出力を読む"),
        ("条件分岐でStatusを承認済",
         "わざと失敗させて、実行履歴の出力から\n原因を見つけ、直します。"),
    ],
    8: [
        ("PA45｜第10回", "PA45｜第11回"),
        ("テーマ解説 (1/2)", "テーマ解説 (1/2)：実行履歴って何？"),
        ("①Forms「新しい応答が送信されるとき」",
         "①「実行履歴」は、フローを動かした記録です。\n　いつ・どのアクションが・成功（緑）か失敗（赤）か が残ります。\n\n②各アクションを開くと「入力」と「出力」が見られます。\n　入力＝受け取った値、出力＝次に渡す値。\n\n③うまく動かないときは、実行履歴の入力・出力を見れば\n　「どこで・何が起きたか」が分かります。"),
    ],
    9: [
        ("PA45｜第10回", "PA45｜第11回"),
        ("テーマ解説 (2/2)", "テーマ解説 (2/2)：画面の「表示」と「値」は違う"),
        ("④「項目の作成」でSharePoint",
         "④画面に見える「表示」と、フローに実際に渡る「値」は\n　違うことがあります。\n\n⑤例：メールの重要度。画面の表示は「高」でも、\n　出力に渡る値は英語の \"High\" です。\n\n⑥条件に「高」と書くと一致せず、フローは赤くならないまま\n　間違った方へ進みます。出力を見ないと気づけません。"),
    ],
    # S10/S11 は Vol.10 事前準備①② を「入力・出力の説明」へ転用
    10: [
        ("PA45｜第10回", "PA45｜第11回"),
        ("事前準備 ①  Formsの設問を作る", "入力と出力って何？"),
        ("🎯 目的：申請者がスマホからも入力できる",
         "🎯 アクションには「入力」と「出力」があります\n\n・入力：そのアクションが受け取った値\n　例）条件アクションの入力 ＝「重要度」と「高」\n\n・出力：そのアクションが次に渡す値\n　例）トリガーの出力 ＝ 届いたメールの件名・差出人・重要度 など\n\n・前のアクションの「出力」が、次のアクションの「入力」に\n　なります。この受け渡しでフローは動いています。"),
    ],
    11: [
        ("PA45｜第10回", "PA45｜第11回"),
        ("事前準備 ②  SharePointリストを作る", "実行履歴で入力・出力を見る"),
        ("🎯 目的：申請データを蓄積して、",
         "🎯 実行履歴を開くと、入力・出力が全部見えます\n\n・フローを実行 →  左メニューの「実行履歴」を開く\n\n・見たい実行をクリック →  各アクションをクリック\n\n・そのアクションの「入力」と「出力」が表示される\n　例）トリガーの出力に、重要度の値 \"High\" が入っている\n\n・思ったとおりに動かないときは、まずここを見る。\n　食い違っている値が、必ずここに写っています。"),
    ],
    12: [
        ("PA45｜第10回", "PA45｜第11回"),
        ("フロー解説 ①  申請を受け取る", "フロー解説 ①  トリガーと条件を置く"),
        ("🎯 目的：紙やメールで来る申請を、",
         "🎯 目的：メールが届いたら重要度で\n　仕分けるフローの形を作る\n\n・トリガー「新しいメールが届いたとき」を置く\n\n・「条件」アクションを追加する\n\n・条件：メールの重要度 が次の値に等しい「高」\n　（まずは、わざと日本語の「高」と書きます）\n\n・はい／いいえ それぞれに「作成」でメッセージを置く"),
    ],
    13: [
        ("PA45｜第10回", "PA45｜第11回"),
        ("フロー解説 ②  登録 → 承認待ち", "フロー解説 ②  わざと失敗させてみる"),
        ("🎯 目的：申請をデータとして残しつつ",
         "🎯 目的：思いどおりに動かない場面を体験する\n\n・フローを保存して、自分宛に「重要度＝高」の\n　メールを1通送ります\n\n・フローは成功（緑）します。でも結果は「通常メール」！\n\n・「重要にして送ったのに、なぜ？」\n　→ ここで実行履歴の出番です"),
    ],
    14: [
        ("PA45｜第10回", "PA45｜第11回"),
        ("フロー解説 ③  結果で更新", "フロー解説 ③  出力で原因を見つけて直す"),
        ("🎯 目的：承認結果をデータに反映し、",
         "🎯 目的：実行履歴の出力から原因を突き止める\n\n・実行履歴を開き、トリガーの「出力」を見る\n\n・重要度（Importance）の値が \"High\"（英語）だと分かる\n\n・条件の「高」を \"High\" に直す\n\n・もう一度、重要度＝高のメールを送る\n　→ 今度こそ「⚠️ 重要メールです」に進む。直った！"),
    ],
    15: [
        ("PA45｜第10回", "PA45｜第11回"),
    ],
    16: [
        ("Vol.10", "Vol.11"),
        ("① Formsは「申請の受け皿」",
         "① フローが成功（緑）でも、結果が正しいとは限らない"),
        ("Formsで入力 → SharePointに保存 の流れは",
         "赤いエラーが出なくても、間違った方へ静かに進むことがあります。\n結果を信じる前に、実行履歴の出力で「本当に合っているか」を確かめましょう。"),
        ("② 承認結果（outcome）は英語固定",
         "② 後続に値を渡すアクションは、まず出力で「本当の値」を確認"),
        ("条件分岐で日本語「承認」と比較すると",
         "条件に書く前に、実行履歴でそのアクションの出力を見る。\n「画面の表示」と「渡る値」は違うことがある、と知っておく。"),
        ("③ 「項目の作成」の {Link}",
         "③ 表示と値は違う ── 同じ罠があちこちに"),
        ("承認カードに項目のリンクを入れておくと",
         "重要度＝High/Normal/Low、承認結果＝Approve/Reject、はい・いいえ列＝true/false。\nどれも英語や真偽値。出力を見れば「本当の値」が分かります。"),
    ],
    17: [
        ("経費申請の自動化", "重要メールを即通知"),
        ("Forms入力→SP登録",
         "重要度＝高のメールが届いたら\nTeamsやスマホへ即お知らせ"),
        ("有休申請の", "差出人で振り分け"),
        ("休暇日・理由をFormsで申請",
         "上司や取引先からのメールだけ\n別処理・別フォルダへ自動仕分け"),
        ("稟議書の", "件名キーワードで仕分け"),
        ("稟議書をFormsで起票",
         "件名に「申請」「至急」を含む\nメールだけ自動で記録・通知"),
        ("実務応用の総まとめ",
         "メールの仕分けは応用がきく"),
    ],
    18: [
        ("PA45｜第10回", "PA45｜第11回"),
    ],
}

# ── 見出しフォント正規化: slide番号 -> [(キー, pt), ...] ──
HEADING_PT = {
    7:  [("今日やること（3つだけ）", 30)],
    8:  [("テーマ解説 (1/2)", 30)],
    9:  [("テーマ解説 (2/2)", 30)],
    10: [("入力と出力って何？", 30)],
    11: [("実行履歴で入力・出力を見る", 30)],
    12: [("フロー解説 ①", 30)],
    13: [("フロー解説 ②", 30)],
    14: [("フロー解説 ③", 30)],
    16: [("今日のポイント", 30)],
    17: [("メールの仕分けは応用がきく", 30)],
}

# ── Vol.10フロー図を除去しスクショ枠を置くスライド ──
# slide番号 -> (枠キャプション)
PLACEHOLDER = {
    8:  "実行履歴の一覧画面\n（緑＝成功・赤＝失敗が見える画面）",
    9:  "メールの重要度を「高」にしている画面\nまたは 条件アクションの設定画面",
    10: "アクションを開いて\n「入力」「出力」が見えている画面",
    11: "実行履歴で トリガーの出力\nが見えている画面",
    12: "完成したフロー全体\n（トリガー＋条件＋作成）",
    13: "実行結果：成功（緑）なのに\n「通常メール」になっている画面",
    14: "トリガーの出力で 重要度が\n\"High\" になっている箇所",
}


def set_text(shape, new_text):
    """shapeのテキストを new_text に置換。先頭runの書式を全体に適用。"""
    tf = shape.text_frame
    lines = new_text.split("\n")
    base = tf.paragraphs[0]
    bsize = base.runs[0].font.size if base.runs else None
    bbold = base.runs[0].font.bold if base.runs else None
    try:
        bname = base.runs[0].font.name if base.runs else None
    except Exception:
        bname = None
    bcolor = None
    if base.runs:
        try:
            if base.runs[0].font.color and base.runs[0].font.color.type is not None:
                bcolor = base.runs[0].font.color.rgb
        except Exception:
            bcolor = None
    # 段落を1つに
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    for r in base.runs[1:]:
        r._r.getparent().remove(r._r)
    if not base.runs:
        base.add_run()
    base.runs[0].text = lines[0]
    for line in lines[1:]:
        np = tf.add_paragraph()
        nr = np.add_run()
        nr.text = line
        if bsize: nr.font.size = bsize
        if bbold is not None: nr.font.bold = bbold
        if bname:
            try: nr.font.name = bname
            except Exception: pass
        if bcolor is not None:
            try: nr.font.color.rgb = bcolor
            except Exception: pass


def main():
    shutil.copy2(SRC, DST)
    prs = Presentation(DST)

    # 1. 文言置換
    n_repl = 0
    for i, slide in enumerate(prs.slides, 1):
        for key, new in REPLACE.get(i, []):
            for shape in slide.shapes:
                if shape.has_text_frame and key in shape.text_frame.text:
                    set_text(shape, new)
                    n_repl += 1
                    break

    # 2. 見出しフォント正規化
    n_head = 0
    for i, items in HEADING_PT.items():
        slide = prs.slides[i - 1]
        for key, pt in items:
            for shape in slide.shapes:
                if shape.has_text_frame and key in shape.text_frame.text:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(pt)
                    n_head += 1
                    break

    # 3. Vol.10フロー図を除去 → スクショ枠を設置
    n_ph = 0
    for i, caption in PLACEHOLDER.items():
        slide = prs.slides[i - 1]
        # 大きい画像（幅2.8in以上 & 高さ3.5in以上）を除去
        for shape in list(slide.shapes):
            if shape.shape_type == 13:  # PICTURE
                try:
                    if shape.width >= Inches(2.8) and shape.height >= Inches(3.5):
                        shape._element.getparent().remove(shape._element)
                except Exception:
                    pass
        # プレースホルダ枠（角丸四角＋キャプション）
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7.7), Inches(1.9), Inches(5.0), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xFB)
        box.line.color.rgb = RGBColor(0x9C, 0xB8, 0xD8)
        box.line.width = Pt(1.5)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = "📷 ここにスクショを貼る"
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = RGBColor(0x2A, 0x5E, 0x9E)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = caption
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0x55, 0x66, 0x77)
        n_ph += 1

    prs.save(DST)
    print(f"文言置換 {n_repl}件 / 見出し正規化 {n_head}件 / スクショ枠 {n_ph}枚")
    print(f"保存: {DST}")


if __name__ == "__main__":
    main()
