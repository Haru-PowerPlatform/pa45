# -*- coding: utf-8 -*-
"""
PA45 Vol.12 PPTX ビルドスクリプト
P011（Vol.11 16スライド）をベースに Vol.12「式アレルギー卒業｜addDays() で3営業日後リマインド」へ書き換え。
- 全テキストの文言を Vol.12 内容へ置換
- スクショ枠のキャプションを Vol.12 用に差し替え
- フォントは Meiryo UI（latin+ea）を全 run に適用
出力: assets/pa45/P012_PA45_FormulaReminder_20260528.pptx

P011 実構造（16スライド）:
 S1 オープニング / S2 事前準備 / S3 タイトル / S4 PA45とは / S5 参加お願い / S6 講師紹介
 S7 今日やること3つ / S8 テーマ解説(1/2) / S9 テーマ解説(2/2)
 S10 入力と出力って何 / S11 実行履歴で見る
 S12 フロー解説① / S13 ハンズオン行きましょう（区切り）
 S14 今日のポイント / S15 活用事例 / S16 アンケート
"""
import sys, io, shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

ROOT = Path(__file__).parent.parent
SRC  = ROOT / "assets" / "pa45" / "P011_PA45_RunHistory_20260521.pptx"
DST  = ROOT / "assets" / "pa45" / "P012_PA45_FormulaReminder_20260528.pptx"
FONT = "Meiryo UI"

# 部分一致キー → 新テキスト全文
REPLACE = {
    1: [
        ("第11回 / Vol.11", "第12回 / Vol.12"),
        ("第11回｜実行履歴を読めると自動化が楽しくなる", "第12回｜「式」アレルギー、今日で卒業"),
        ("失敗しないフロー設計｜メールの重要度で振り分け",
         "申請から3営業日後リマインド｜addDays() 1個で書ける"),
    ],
    3: [
        ("Vol.11", "Vol.12"),
        ("PA45【メール受信", "PA45【式 addDays × 申請 × リマインド】"),
        ("失敗しないフロー設計\n実行履歴を読めるようになろう",
         "第12回： \n「式」アレルギー卒業\naddDays() で日付計算を体験"),
        ("Theme　メールの重要度",
         "Theme　申請から3営業日後リマインド｜式 addDays() を初体験"),
        ("Target　条件分岐（Vol.3）",
         "Target　式に触れたことがない方／式を復習したい方も歓迎"),
    ],
    7: [
        ("PA45｜第11回", "PA45｜第12回"),
        ("メール受信でフローを作る", "申請日を入力するフローを作る"),
        ("「新しいメールが届いたとき」をトリガーに",
         "手動トリガーで「申請日」を入力。\nPower Automate に日付を渡すところから始めます。"),
        ("重要度で条件分岐させる", "addDays() で3日後を計算する"),
        ("「条件」アクションで",
         "「式」タブから addDays() を使って、\n申請日から3日後を出します。"),
        ("実行履歴の出力を読む", "リマインド文を作って通知"),
        ("わざと失敗させて",
         "計算した日付をリマインド文に組み込んで、\n自分宛にメールやTeamsで通知します。"),
    ],
    8: [
        ("PA45｜第11回", "PA45｜第12回"),
        ("テーマ解説 (1/2)", "テーマ解説 (1/2)：「式」って怖くない"),
        ("①「実行履歴」は、フローを動かした記録",
         "①「式」はちょっとした計算ルールです。動的コンテンツの隣にある「fx 式」タブから書けます。"),
        ("②各アクションを開くと",
         "②難しそうな見た目ですが、Excel の関数を1個書くのと同じ感覚です。"),
        ("③うまく動かないときは",
         "③関数を1個覚えるだけで、できることが一気に広がります。今日は addDays() ひとつだけ。"),
    ],
    9: [
        ("PA45｜第11回", "PA45｜第12回"),
        ("テーマ解説 (2/2)", "テーマ解説 (2/2)：addDays() の使い方"),
        ("④画面に見える「表示」",
         "④addDays(基準日, 加算日数, 書式) という3つを書きます。"),
        ("⑤例：メールの重要度",
         "⑤例：addDays('2026-05-28', 3, 'yyyy/MM/dd') と書くと「2026-05-31」が返ります。"),
        ("⑥条件に「高」と書くと",
         "⑥動的コンテンツの「申請日」を基準日に渡せば、申請日からN日後を自動で計算できます。"),
    ],
    10: [
        ("PA45｜第11回", "PA45｜第12回"),
        ("入力と出力って何？", "「式」を書く場所"),
        ("🎯 アクションには「入力」",
         "🎯 式は「動的コンテンツ」の隣の「fx」タブから書きます"),
        ("・入力：そのアクションが受け取った値",
         "・アクションの入力欄をクリック → 右側にパネルが開く"),
        ("・出力：そのアクションが次に渡す値",
         "・パネルの上に「動的なコンテンツ」と「式」の2つのタブがある"),
        ("・前のアクションの「出力」が",
         "・「式」タブを選ぶと、関数を入力する欄が出る → 関数名を入れると候補が出て補完できる"),
    ],
    11: [
        ("PA45｜第11回", "PA45｜第12回"),
        ("実行履歴で入力・出力を見る", "addDays() を実際に書いてみる"),
        ("🎯 実行履歴を開くと",
         "🎯 動的コンテンツと関数を組み合わせます"),
        ("・フローを実行 → 左メニューの",
         "・「式」タブで `addDays(` まで入力"),
        ("・見たい実行をクリック",
         "・「動的なコンテンツ」タブに切り替え → 「申請日」をクリック"),
        ("・そのアクションの「入力」と「出力」",
         "・式タブに戻って `, 3, 'yyyy/MM/dd')` を追加してOKを押す"),
        ("・思ったとおりに動かないとき",
         "・完成形：addDays(triggerBody()?['申請日'], 3, 'yyyy/MM/dd')"),
    ],
    12: [
        ("PA45｜第11回", "PA45｜第12回"),
        ("フロー解説 ①  トリガーと条件を置く", "フロー解説 ①  申請日 → addDays → 通知"),
        ("🎯 目的：メールが届いたら",
         "🎯 目的：申請日を受け取って、3日後にリマインドを通知する"),
        ("・トリガー「新しいメールが届いたとき」を置く",
         "・トリガー「手動でフローをトリガーします」＋入力「申請日」（テキスト）"),
        ("・「条件」アクションを追加する",
         "・「作成」アクションで式 addDays(申請日, 3, 'yyyy/MM/dd') を書く"),
        ("・条件：メールの重要度",
         "・もう1つ「作成」でリマインド文を組む（「3日後の◯◯にリマインドです」）"),
        ("・はい／いいえ それぞれに",
         "・最後に「メールを送信」または「Teams で投稿」で自分宛に通知"),
    ],
    13: [
        # ハンズオン行きましょう（区切り）— ラベルだけ更新
        ("PA45｜第11回", "PA45｜第12回"),
    ],
    14: [
        # 今日のポイント
        ("Vol.11", "Vol.12"),
        ("① フローが成功（緑）でも、結果が正しいとは限らない",
         "① 「式」は怖くない、動的コンテンツの「隣」にある"),
        ("赤いエラーが出なくても",
         "「式」タブは fx ボタンを押すだけで開きます。書くのは関数1個から。"),
        ("② 後続に値を渡すアクションは、まず出力で",
         "② addDays() 1個覚えれば、日付計算は怖くない"),
        ("条件に書く前に、実行履歴",
         "addDays(基準日, 加算日数, 書式) のパターンさえ覚えれば、応用が一気に広がります。"),
        ("③ 表示と値は違う ── 同じ罠があちこちに",
         "③ 困ったときは、実行履歴で「式の結果」を確認"),
        ("重要度＝High/Normal/Low",
         "出力を開けば、式が計算した結果が見えます。式が思いどおりかどうかは、出力でチェック。"),
    ],
    15: [
        # 活用事例
        ("メールの仕分けは応用がきく", "日付の計算は応用がきく"),
    ],
    16: [
        ("PA45｜第11回", "PA45｜第12回"),
    ],
}

HEADING_PT = {
    7:  [("今日やること（3つだけ）", 30)],
    8:  [("テーマ解説 (1/2)", 30)],
    9:  [("テーマ解説 (2/2)", 30)],
    10: [("「式」を書く場所", 30)],
    11: [("addDays() を実際に書いてみる", 30)],
    12: [("フロー解説 ①", 30)],
    14: [("今日のポイント", 30)],
    15: [("日付の計算は応用がきく", 30)],
}

PLACEHOLDER_CAPTIONS = {
    8:  "「式」タブを開いた画面\n（fx ボタンの場所がわかる画面）",
    9:  "addDays() の式を書いている画面\n（補完候補が出ている状態）",
    10: "アクションの入力欄の右側\n「動的なコンテンツ／式」タブ切り替え",
    11: "addDays(動的コンテンツ, 3, 書式)\nを書き終えた完成画面",
    12: "完成フロー全体\n（トリガー→addDays→リマインド送信）",
}

# 活用事例（S15）の各カードを index 指定で置換（重複テキスト対策）
S15_INDEX_REPLACE = {
    7:  "提出期限リマインド",
    8:  "申請日から3日後を\n自動で通知",
    9:  "申請日を受け取る\n→ addDays(+3日)\n→ Teamsで通知",
    11: "申請日 → addDays → 通知",
    16: "契約更新前の通知",
    17: "契約日から1年後の30日前に\n更新検討のリマインド",
    18: "契約日を受け取る\n→ addDays(+335日)\n→ メールで通知",
    20: "契約日 → addDays → 通知",
    25: "受講期限リマインド",
    26: "研修や健康診断の受講期限が\n近づいたら自動で通知",
    27: "受講日を受け取る\n→ addDays(-7日)\n→ リマインドメール",
}

def set_font_name(run):
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea'):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set('typeface', FONT)

def set_text(shape, new_text):
    tf = shape.text_frame
    lines = new_text.split("\n")
    base = tf.paragraphs[0]
    bsize = base.runs[0].font.size if base.runs else None
    bbold = base.runs[0].font.bold if base.runs else None
    bcolor = None
    if base.runs:
        try:
            if base.runs[0].font.color and base.runs[0].font.color.type is not None:
                bcolor = base.runs[0].font.color.rgb
        except Exception:
            pass
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    for r in base.runs[1:]:
        r._r.getparent().remove(r._r)
    if not base.runs:
        base.add_run()
    base.runs[0].text = lines[0]
    if bsize: base.runs[0].font.size = bsize
    if bbold is not None: base.runs[0].font.bold = bbold
    if bcolor is not None:
        try: base.runs[0].font.color.rgb = bcolor
        except: pass
    set_font_name(base.runs[0])
    if lines[0] and (lines[0][0] in '①②③④⑤⑥' or lines[0].startswith('🎯')):
        base.space_before = Pt(10)
    for line in lines[1:]:
        np = tf.add_paragraph()
        nr = np.add_run(); nr.text = line
        if bsize: nr.font.size = bsize
        if bbold is not None: nr.font.bold = bbold
        if bcolor is not None:
            try: nr.font.color.rgb = bcolor
            except: pass
        set_font_name(nr)
        if line and (line[0] in '①②③④⑤⑥' or line.startswith('🎯')):
            np.space_before = Pt(10)

def main():
    shutil.copy2(SRC, DST)
    prs = Presentation(DST)

    # 1. 文言置換
    n_repl = 0
    for i, slide in enumerate(prs.slides, 1):
        for key, new in REPLACE.get(i, []):
            for shape in slide.shapes:
                if shape.has_text_frame and key in shape.text_frame.text:
                    set_text(shape, new)
                    n_repl += 1
                    break

    # 2. S15（活用事例）の index 指定置換
    s15 = list(prs.slides[14].shapes)
    n_s15 = 0
    for idx, txt in S15_INDEX_REPLACE.items():
        if idx < len(s15):
            set_text(s15[idx], txt)
            n_s15 += 1
    print(f'S15 カード関連: {n_s15}箇所更新')

    # 3. 見出しフォント正規化
    n_head = 0
    for i, items in HEADING_PT.items():
        slide = prs.slides[i - 1]
        for key, pt in items:
            for shape in slide.shapes:
                if shape.has_text_frame and key in shape.text_frame.text:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(pt)
                    n_head += 1
                    break

    # 4. スクショ枠キャプション差し替え
    n_ph = 0
    for i, caption in PLACEHOLDER_CAPTIONS.items():
        slide = prs.slides[i - 1]
        for shape in slide.shapes:
            if shape.has_text_frame and "ここにスクショを貼る" in shape.text_frame.text:
                tf = shape.text_frame
                paras = list(tf.paragraphs)
                if len(paras) >= 2:
                    p2 = paras[1]
                    bsize = p2.runs[0].font.size if p2.runs else Pt(14)
                    for r in p2.runs:
                        r._r.getparent().remove(r._r)
                    lines = caption.split('\n')
                    nr = p2.add_run(); nr.text = lines[0]
                    nr.font.size = bsize; set_font_name(nr)
                    for extra in lines[1:]:
                        np = tf.add_paragraph()
                        nrr = np.add_run(); nrr.text = extra
                        nrr.font.size = bsize; set_font_name(nrr)
                    # 既存の3段目以降を削除
                    keep_n = 1 + len(lines)
                    for excess in list(tf.paragraphs)[keep_n:]:
                        excess._p.getparent().remove(excess._p)
                n_ph += 1
                break

    # 5. 全 run フォント Meiryo UI 統一
    n_font = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    set_font_name(run); n_font += 1

    prs.save(DST)
    print(f'文言置換 {n_repl}件 / 見出し {n_head}枚 / スクショ枠 {n_ph}枚 / フォント {n_font} run')
    print(f'保存: {DST}')

if __name__ == "__main__":
    main()
